"""文档摄入层 — 替代 file-parser.js

支持 64+ 格式，使用 Unstructured.io (主) + MarkItDown (轻量降级)

Pipeline:
  Unstructured.io (hi_res) → MarkItDown → python-docx/PyMuPDF/pdfplumber → 纯文本
"""
import io
import os
import logging
from pathlib import Path
from typing import Optional

from backend.config import UNSTRUCTURED_ENABLED, UNSTRUCTURED_STRATEGY

logger = logging.getLogger(__name__)

# 动态导入，避免缺少依赖时整个模块加载失败
_unstructured_available = False
_markitdown_available = False
_pymupdf_available = False

try:
    from unstructured.partition.auto import partition
    from unstructured.staging.base import elements_to_json
    _unstructured_available = True
except ImportError:
    logger.info("unstructured not installed — using fallback parsers")

try:
    from markitdown import MarkItDown
    _markitdown_available = True
    _md_converter = MarkItDown()
except ImportError:
    logger.info("markitdown not installed — using basic parsers")

try:
    import fitz  # PyMuPDF
    _pymupdf_available = True
except ImportError:
    logger.info("PyMuPDF not installed — PDF text extraction limited")


class IngestedDocument:
    """摄入后的文档"""
    def __init__(self):
        self.markdown: str = ""
        self.raw_text: str = ""
        self.doc_type: str = "unknown"
        self.structure: Optional[dict] = None  # 版面条目、表格等
        self.method: str = "unknown"
        self.page_count: int = 0


class DocumentIngestion:
    """文档摄入层"""

    # 支持的 MIME 类型映射
    MIME_MAP = {
        "application/pdf": "pdf",
        "image/png": "image",
        "image/jpeg": "image",
        "image/jpg": "image",
        "image/webp": "image",
        "image/bmp": "image",
        "image/gif": "image",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "docx",
        "application/msword": "doc",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": "xlsx",
        "application/vnd.ms-excel": "xls",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx",
        "text/plain": "txt",
        "text/html": "html",
        "text/csv": "csv",
    }

    async def ingest(self, file_bytes: bytes, file_name: str, mime_type: str = "") -> IngestedDocument:
        """主入口：摄入任意文档，返回结构化文本"""
        doc = IngestedDocument()
        doc_type = self._detect_type(file_name, mime_type)
        doc.doc_type = doc_type

        # Level 1: Unstructured.io
        if UNSTRUCTURED_ENABLED and _unstructured_available:
            try:
                return await self._ingest_unstructured(file_bytes, file_name, doc_type, doc)
            except Exception as e:
                logger.warning(f"Unstructured failed: {e}, falling back")

        # Level 2: MarkItDown
        if _markitdown_available:
            try:
                return await self._ingest_markitdown(file_bytes, file_name, doc_type, doc)
            except Exception as e:
                logger.warning(f"MarkItDown failed: {e}, falling back")

        # Level 3: Type-specific parsers
        return await self._ingest_fallback(file_bytes, file_name, doc_type, doc)

    async def _ingest_unstructured(self, file_bytes: bytes, file_name: str, doc_type: str, doc: IngestedDocument) -> IngestedDocument:
        """使用 Unstructured.io 解析"""
        file_io = io.BytesIO(file_bytes)

        elements = partition(
            file=file_io,
            content_type=self._get_content_type(doc_type),
            strategy=UNSTRUCTURED_STRATEGY,
        )

        # 按类型分类元素
        texts = []
        tables = []
        for el in elements:
            el_dict = el.to_dict()
            if el_dict.get("type") == "Table":
                tables.append(el_dict)
            texts.append(str(el))

        doc.markdown = "\n\n".join(texts)
        doc.raw_text = "\n".join(texts)
        doc.structure = {"tables": tables, "element_count": len(elements)}
        doc.method = "unstructured"

        logger.info(f"[Unstructured] Parsed {doc_type}: {len(elements)} elements, {len(tables)} tables")
        return doc

    async def _ingest_markitdown(self, file_bytes: bytes, file_name: str, doc_type: str, doc: IngestedDocument) -> IngestedDocument:
        """使用 MarkItDown 转换"""
        # 保存临时文件（MarkItDown 需要文件路径）
        import tempfile
        suffix = f".{doc_type}" if doc_type != "unknown" else ""
        with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name

        try:
            result = _md_converter.convert(tmp_path)
            doc.markdown = result.text_content
            doc.raw_text = result.text_content
            doc.method = "markitdown"
        finally:
            os.unlink(tmp_path)

        logger.info(f"[MarkItDown] Converted {doc_type}: {len(doc.markdown)} chars")
        return doc

    async def _ingest_fallback(self, file_bytes: bytes, file_name: str, doc_type: str, doc: IngestedDocument) -> IngestedDocument:
        """类型特定的降级解析器"""
        if doc_type == "pdf":
            return self._parse_pdf_fallback(file_bytes, doc)
        elif doc_type in ("docx", "doc"):
            return self._parse_docx_fallback(file_bytes, doc)
        elif doc_type in ("xlsx", "xls"):
            return self._parse_xlsx_fallback(file_bytes, doc)
        elif doc_type in ("pptx", "ppt"):
            return self._parse_pptx_fallback(file_bytes, doc)
        elif doc_type == "image":
            return self._parse_image_fallback(file_bytes, doc)
        else:
            # 默认：尝试作为文本读取
            try:
                text = file_bytes.decode("utf-8")
                doc.markdown = text
                doc.raw_text = text
                doc.method = "plain_text"
            except UnicodeDecodeError:
                doc.markdown = f"[无法解析二进制文件: {doc_type}]"
                doc.method = "failed"
        return doc

    def _parse_pdf_fallback(self, file_bytes: bytes, doc: IngestedDocument) -> IngestedDocument:
        """PyMuPDF 降级 PDF 解析"""
        if _pymupdf_available:
            pdf = fitz.open(stream=file_bytes, filetype="pdf")
            doc.page_count = pdf.page_count
            texts = []
            for page in pdf:
                texts.append(page.get_text())
            text = "\n\n".join(texts)
            doc.markdown = text
            doc.raw_text = text
            doc.method = "pymupdf"
            logger.info(f"[PyMuPDF] Parsed {doc.page_count} pages")
        else:
            # 最后的兜底：简单文本提取
            text = file_bytes.decode("utf-8", errors="ignore")
            # 尝试提取可读文本
            import re
            # 移除 PDF 二进制控制字符
            text = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f]', '', text)
            doc.markdown = text.strip()
            doc.raw_text = doc.markdown
            doc.method = "plain_text_fallback"
        return doc

    def _parse_docx_fallback(self, file_bytes: bytes, doc: IngestedDocument) -> IngestedDocument:
        """python-docx 降级 DOCX 解析"""
        try:
            from docx import Document
            from io import BytesIO
            d = Document(BytesIO(file_bytes))
            texts = [p.text for p in d.paragraphs if p.text.strip()]
            text = "\n\n".join(texts)
            doc.markdown = text
            doc.raw_text = text
            doc.method = "python-docx"
        except Exception:
            text = file_bytes.decode("utf-8", errors="ignore")
            # 提取 <w:t> 标签文本
            import re
            t_tags = re.findall(r'<w:t[^>]*>(.*?)</w:t>', text)
            doc.markdown = "".join(t_tags)
            doc.raw_text = doc.markdown
            doc.method = "docx_xml_fallback"
        return doc

    def _parse_xlsx_fallback(self, file_bytes: bytes, doc: IngestedDocument) -> IngestedDocument:
        """openpyxl 降级 Excel 解析"""
        try:
            from openpyxl import load_workbook
            from io import BytesIO
            wb = load_workbook(BytesIO(file_bytes), read_only=True, data_only=True)
            texts = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                texts.append(f"## {sheet_name}")
                for row in ws.iter_rows(values_only=True):
                    row_text = " | ".join(str(c) if c is not None else "" for c in row)
                    if row_text.strip():
                        texts.append(row_text)
            text = "\n".join(texts)
            doc.markdown = text
            doc.raw_text = text
            doc.method = "openpyxl"
        except Exception as e:
            doc.markdown = f"[Excel 解析失败: {e}]"
            doc.method = "failed"
        return doc

    def _parse_pptx_fallback(self, file_bytes: bytes, doc: IngestedDocument) -> IngestedDocument:
        """python-pptx 降级 PPTX 解析"""
        try:
            from pptx import Presentation
            from io import BytesIO
            prs = Presentation(BytesIO(file_bytes))
            texts = []
            for i, slide in enumerate(prs.slides, 1):
                slide_texts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for p in shape.text_frame.paragraphs:
                            if p.text.strip():
                                slide_texts.append(p.text)
                if slide_texts:
                    texts.append(f"## Slide {i}\n" + "\n".join(slide_texts))
            text = "\n\n".join(texts)
            doc.markdown = text
            doc.raw_text = text
            doc.method = "python-pptx"
        except Exception as e:
            doc.markdown = f"[PPTX 解析失败: {e}]"
            doc.method = "failed"
        return doc

    def _parse_image_fallback(self, file_bytes: bytes, doc: IngestedDocument) -> IngestedDocument:
        """图片降级处理 — 通常需要 OCR，此处标记"""
        doc.markdown = "[图片文件 — 需要 OCR 处理]"
        doc.raw_text = ""
        doc.method = "image_pending_ocr"
        return doc

    def _detect_type(self, file_name: str, mime_type: str) -> str:
        """检测文件类型"""
        # 优先使用 MIME 类型
        if mime_type in self.MIME_MAP:
            return self.MIME_MAP[mime_type]

        # 通过扩展名
        ext = Path(file_name).suffix.lower().lstrip(".")
        ext_map = {
            "pdf": "pdf", "png": "image", "jpg": "image", "jpeg": "image",
            "webp": "image", "bmp": "image", "gif": "image",
            "docx": "docx", "doc": "doc", "xlsx": "xlsx", "xls": "xls",
            "pptx": "pptx", "ppt": "ppt", "txt": "txt", "html": "html",
            "csv": "csv", "md": "txt", "json": "txt",
        }
        return ext_map.get(ext, "unknown")

    def _get_content_type(self, doc_type: str) -> str:
        """获取 Unstructured 的内容类型"""
        type_map = {
            "pdf": "application/pdf",
            "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        }
        return type_map.get(doc_type, "text/plain")
